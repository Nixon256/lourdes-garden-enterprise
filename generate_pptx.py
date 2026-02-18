from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define brand colors
    FOREST_GREEN = RGBColor(5, 46, 22) # #052e16
    GOLD = RGBColor(22, 163, 74) # Using green-600 logic for consistency

    def add_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Style Slide
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = FOREST_GREEN
        title.text_frame.paragraphs[0].font.bold = True

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = bullet_points[0]

        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.space_after = Pt(10)

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Lourdes Garden V1.0"
    subtitle.text = "Premium Global Agricultural Heritage | Hitting 100/100 Perfection"

    # Slide 2: Mission & Brand
    add_slide("Mission & Brand Identity", [
        "Brand Vision: 'From our heritage grove in Tamil Nadu to the global stage.'",
        "Aesthetic: 'Quiet Luxury' & Luxe-Editorial Cinematic Storytelling.",
        "Focus: Sustainable Organic Farming, Export Quality, Botanical Heritage.",
        "Vibe: Minimalist, sophisticated, and professional."
    ])

    # Slide 3: Technical Stack
    add_slide("Technical Architecture", [
        "Core: Next.js 16 (App Router) & React 19.",
        "Type Safety: 100% TypeScript (End-to-End).",
        "Database: Serverless PostgreSQL (Neon) & Prisma ORM.",
        "Performance: Framer Motion interactions & Tailwind CSS.",
        "Infrastructure: Vercel Cloud Deployment Ready."
    ])

    # Slide 4: Feature Highlights
    add_slide("V1.0 Feature Highlights", [
        "Full Bilingual Experience: English & Tamil native toggle support.",
        "Narrative Products: Cinematic product articles focusing on soul and soil.",
        "Art of the Soil: 59+ optimized heritage assets in masonry gallery.",
        "Contact Portal: Secure, rate-limited inquiry microservice with DB persistence."
    ])

    # Slide 5: Quality Perfection (The 100/100 Score)
    add_slide("Quality & Launch Readiness", [
        "SEO: 10/10 (Dynamic Sitemap, Robots, JSON-LD Organization Schema).",
        "Aesthetic: 10/10 (Bespoke Botanical Cursor, Arima Typography).",
        "Performance: 10/10 (PWA Manifest, Optimized Pre-fetching).",
        "Status: Platinum Ready for Public Launch."
    ])

    # Slide 6: Page Map
    add_slide("Sitemap & User Journey", [
        "Home: Hero discovery and 'Featured Harvest' highlights.",
        "About: Foundation history (Since 2020) and Mission.",
        "Products: Cinematic journey through mountain-grown organics.",
        "Gallery: High-fidelity visual evidence of excellence.",
        "Contact: Global reach portal and farm location."
    ])

    # Slide 7: Future (V2.0)
    add_slide("Future Roadmap (V2.0)", [
        "B2B Enterprise Dashboard & Order Tracking.",
        "Direct Global Checkout with multi-currency support.",
        "Real-time Farm-to-Table transparency engine."
    ])

    output_path = "Lourdes_Garden_V1_Excellence_Presentation.pptx"
    prs.save(output_path)
    print(f"PPT_CREATED_AT:{output_path}")

if __name__ == "__main__":
    create_presentation()
